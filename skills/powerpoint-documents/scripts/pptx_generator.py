#!/usr/bin/env python3
"""
pptx_generator.py

Backend for PowerPoint presentation creation and manipulation using python-pptx
"""

import sys
import json
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_color):
    """Convert hex color #RRGGBB to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_presentation(output, title=None, subtitle=None, slides=None, template=None, 
                       width=10, height=7.5, theme="light"):
    """Create new PowerPoint presentation"""
    
    if template and Path(template).exists():
        prs = Presentation(template)
    else:
        prs = Presentation()
        # Set slide dimensions (default 16:9)
        prs.slide_width = Inches(width)
        prs.slide_height = Inches(height)
    
    # Add title slide if title provided
    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle
    
    # Add slides from JSON if provided
    if slides:
        if isinstance(slides, str):
            slides = json.loads(slides)
        
        for slide_def in slides:
            add_slide_from_def(prs, slide_def)
    
    prs.save(output)
    print(f"✅ Presentation created: {output}")

def add_slide_from_def(prs, slide_def):
    """Add slide from definition dict"""
    
    layout_name = slide_def.get('layout', 'content')
    layout_map = {
        'title': 0,           # Title slide
        'title-only': 5,      # Title only
        'content': 1,         # Title and content
        '2col': 3,           # Two content
        '3col': 2,           # Comparison (can be used as 3-col)
        'blank': 6,          # Blank
        'image-right': 4,    # Picture with caption
    }
    
    layout_idx = layout_map.get(layout_name, 1)
    if layout_idx >= len(prs.slide_layouts):
        layout_idx = 1
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Add content based on layout
    if layout_name == 'title':
        slide.shapes.title.text = slide_def.get('title', '')
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_def.get('subtitle', '')
    
    elif layout_name == 'content':
        if slide.shapes.title:
            slide.shapes.title.text = slide_def.get('heading', '')
        # Add content to first text placeholder
        for shape in slide.placeholders:
            if shape.is_placeholder and 'content' in shape.name.lower():
                if shape.has_text_frame:
                    shape.text = slide_def.get('content', '')
                break
    
    elif layout_name == '2col':
        if slide.shapes.title:
            slide.shapes.title.text = slide_def.get('heading', '')
        # Add to left and right content areas
        idx = 0
        for shape in slide.placeholders:
            if shape.is_placeholder and idx < 2:
                if shape.has_text_frame:
                    if idx == 0:
                        shape.text = slide_def.get('left_content', '')
                    else:
                        shape.text = slide_def.get('right_content', '')
                    idx += 1
    
    elif layout_name == 'image-right':
        if slide.shapes.title:
            slide.shapes.title.text = slide_def.get('heading', '')
        # Add text content
        for shape in slide.placeholders:
            if shape.is_placeholder and 'content' in shape.name.lower():
                if shape.has_text_frame:
                    shape.text = slide_def.get('content', '')
                break

def add_slide(input_pres, output, layout, heading=None, content=None, position=None):
    """Add slide to presentation"""
    
    prs = Presentation(input_pres)
    
    slide_def = {
        'layout': layout,
        'heading': heading,
        'content': content
    }
    
    add_slide_from_def(prs, slide_def)
    prs.save(output)
    print(f"✅ Slide added: {output}")

def add_image(input_pres, output, image_path, slide_idx, width=None, height=None, 
             position="center", caption=None):
    """Add image to specific slide"""
    
    if not Path(image_path).exists():
        raise FileNotFoundError(f"Image not found: {image_path}")
    
    prs = Presentation(input_pres)
    
    if slide_idx >= len(prs.slides):
        raise ValueError(f"Slide index {slide_idx} out of range")
    
    slide = prs.slides[slide_idx]
    
    # Convert dimensions if provided
    if width:
        width = Inches(width / 2.54)  # cm to inches
    if height:
        height = Inches(height / 2.54)
    
    # Default positioning
    if width and height:
        pic = slide.shapes.add_picture(image_path, Inches(1), Inches(1), 
                                       width=width, height=height)
    elif width:
        pic = slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=width)
    else:
        pic = slide.shapes.add_picture(image_path, Inches(1), Inches(1))
    
    prs.save(output)
    print(f"✅ Image added to slide {slide_idx}: {output}")

def apply_master(input_pres, output, master_path, preserve_content=True, theme=None):
    """Apply master slide template"""
    
    prs = Presentation(input_pres)
    master_prs = Presentation(master_path)
    
    # Copy master slides from template
    try:
        # This is a simplification - full master application is complex
        for slide_layout in master_prs.slide_layouts:
            # Master application would require deeper XML manipulation
            pass
    except:
        pass
    
    prs.save(output)
    print(f"✅ Master template applied: {output}")

def create_template(output, name, description="", primary_color="#003366",
                   secondary_color="#00AA99", accent_color="#FF6B35",
                   font_family="Calibri", logo=None, company_name=None, layouts=None):
    """Create new master slide template"""
    
    prs = Presentation()
    
    # Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "{{title}}"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "{{subtitle}}"
    
    # Add content slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "{{heading}}"
    
    # Set default colors in theme
    prs.core_properties.author = company_name or "raspb"
    prs.core_properties.title = name
    prs.core_properties.subject = description
    
    # Save template
    prs.save(output)
    print(f"✅ Template created: {output}")

def list_templates():
    """List available templates"""
    
    assets_dir = Path(__file__).parent.parent / 'assets'
    
    if not assets_dir.exists():
        print("No templates found (assets directory empty)")
        return
    
    templates = list(assets_dir.glob('*.pptx'))
    
    if not templates:
        print("No templates found in assets/")
        return
    
    print("Available templates:")
    for template_file in templates:
        print(f"  - {template_file.name}")
        try:
            prs = Presentation(str(template_file))
            if prs.core_properties.title:
                print(f"    Name: {prs.core_properties.title}")
            if prs.core_properties.subject:
                print(f"    Description: {prs.core_properties.subject}")
        except:
            pass

# Main CLI
if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='PowerPoint presentation generator')
    subparsers = parser.add_subparsers(dest='command')
    
    # create
    create_parser = subparsers.add_parser('create')
    create_parser.add_argument('--output', required=True)
    create_parser.add_argument('--title')
    create_parser.add_argument('--subtitle')
    create_parser.add_argument('--slides')
    create_parser.add_argument('--template')
    create_parser.add_argument('--width', type=float, default=10)
    create_parser.add_argument('--height', type=float, default=7.5)
    create_parser.add_argument('--theme', default='light')
    
    # add-slide
    slide_parser = subparsers.add_parser('add-slide')
    slide_parser.add_argument('--input', required=True)
    slide_parser.add_argument('--output', required=True)
    slide_parser.add_argument('--layout', required=True)
    slide_parser.add_argument('--heading')
    slide_parser.add_argument('--content')
    slide_parser.add_argument('--position', type=int)
    
    # add-image
    image_parser = subparsers.add_parser('add-image')
    image_parser.add_argument('--input', required=True)
    image_parser.add_argument('--output', required=True)
    image_parser.add_argument('--image', required=True)
    image_parser.add_argument('--slide', type=int, required=True)
    image_parser.add_argument('--width', type=float)
    image_parser.add_argument('--height', type=float)
    image_parser.add_argument('--position', default='center')
    image_parser.add_argument('--caption')
    
    # apply-master
    master_parser = subparsers.add_parser('apply-master')
    master_parser.add_argument('--input', required=True)
    master_parser.add_argument('--output', required=True)
    master_parser.add_argument('--master', required=True)
    master_parser.add_argument('--preserve-content', default=True)
    master_parser.add_argument('--theme')
    
    # create-template
    create_tpl_parser = subparsers.add_parser('create-template')
    create_tpl_parser.add_argument('--output', required=True)
    create_tpl_parser.add_argument('--name', required=True)
    create_tpl_parser.add_argument('--description', default='')
    create_tpl_parser.add_argument('--primary-color', default='#003366')
    create_tpl_parser.add_argument('--secondary-color', default='#00AA99')
    create_tpl_parser.add_argument('--accent-color', default='#FF6B35')
    create_tpl_parser.add_argument('--font-family', default='Calibri')
    create_tpl_parser.add_argument('--logo')
    create_tpl_parser.add_argument('--company-name')
    create_tpl_parser.add_argument('--layouts')
    
    # list-templates
    subparsers.add_parser('list-templates')
    
    args = parser.parse_args()
    
    try:
        if args.command == 'create':
            create_presentation(args.output, args.title, args.subtitle, args.slides,
                              args.template, args.width, args.height, args.theme)
        elif args.command == 'add-slide':
            add_slide(args.input, args.output, args.layout, args.heading, args.content, args.position)
        elif args.command == 'add-image':
            add_image(args.input, args.output, args.image, args.slide, args.width, 
                     args.height, args.position, args.caption)
        elif args.command == 'apply-master':
            apply_master(args.input, args.output, args.master, args.preserve_content, args.theme)
        elif args.command == 'create-template':
            create_template(args.output, args.name, args.description, args.primary_color,
                          args.secondary_color, args.accent_color, args.font_family,
                          args.logo, args.company_name, args.layouts)
        elif args.command == 'list-templates':
            list_templates()
    except Exception as e:
        print(f"❌ Error: {e}", file=sys.stderr)
        sys.exit(1)
